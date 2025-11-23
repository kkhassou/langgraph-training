"""IO Loader Node - 汎用ファイル読み込みノード

対応形式:
- PowerPoint (.pptx)
- (将来) PDF (.pdf)
- (将来) Word (.docx)
- (将来) Text (.txt, .md)
"""

import os
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation

from src.nodes.base import BaseNode, NodeState, NodeInput, NodeOutput


class LoaderNode(BaseNode):
    """汎用ファイル読み込みノード
    
    拡張子に基づいて自動的に適切なローダーを選択します。
    
    State入力:
        - data["file_path"]: ファイルパス（必須）
    
    State出力:
        - data["content"]: 読み込まれた内容（形式はファイルタイプによる）
        - data["metadata"]: ファイルメタデータ
    """

    def __init__(self):
        super().__init__(
            name="loader_node",
            description="Load content from various file formats"
        )

    async def execute(self, state: NodeState) -> NodeState:
        """ファイル読み込みを実行"""
        try:
            file_path = state.data.get("file_path")
            if not file_path:
                raise ValueError("file_path is required in state.data")

            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            # 拡張子に基づいてロード
            ext = path.suffix.lower()
            content = None

            if ext == ".pptx":
                content = await self._load_pptx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")

            # 結果を格納
            state.data["content"] = content
            state.data["metadata"] = {
                "file_name": path.name,
                "file_size": path.stat().st_size,
                "file_type": ext
            }
            
            # 後方互換性のために従来のキーも維持（必要に応じて削除）
            if ext == ".pptx":
                state.data["extracted_text"] = content
                state.data["slide_count"] = len(content)

            state.messages.append(f"Loaded {ext} file: {path.name}")
            state.metadata["node"] = self.name

            return state

        except Exception as e:
            state.data["error"] = str(e)
            state.metadata["error_node"] = self.name
            return state

    async def _load_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """PPTXファイルを読み込み"""
        slides_text = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "notes": ""
            }

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    
                    # Try to identify title
                    is_title = False
                    try:
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                is_title = True
                    except Exception:
                        pass
                    
                    if is_title:
                        slide_content["title"] = text_content
                    else:
                        slide_content["content"].append(text_content)

            # Extract notes
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide_content["notes"] = slide.notes_slide.notes_text_frame.text.strip()

            slides_text.append(slide_content)

        return slides_text


# ✅ 後方互換性のためのエイリアス
class PowerPointIngestNode(LoaderNode):
    """PowerPoint Ingest Node (deprecated, use LoaderNode)"""
    pass


class LoaderInput(NodeInput):
    """Input model for Loader node"""
    file_path: str


class LoaderOutput(NodeOutput):
    """Output model for Loader node"""
    content: Any = None
    metadata: Dict[str, Any] = {}
    
    # 後方互換性フィールド
    slide_count: int = 0
    extracted_slides: List[Dict[str, Any]] = []


# ✅ 後方互換性のためのハンドラー
async def ppt_ingest_handler(file_path: str) -> LoaderOutput:
    """Standalone handler for PowerPoint ingest API endpoint"""
    try:
        node = LoaderNode()
        state = NodeState()
        state.data["file_path"] = file_path

        result_state = await node.execute(state)

        if "error" in result_state.data:
            return LoaderOutput(
                output_text="",
                success=False,
                error_message=result_state.data["error"]
            )

        content = result_state.data.get("content", [])
        
        # Format output text for PPT
        output_lines = []
        for slide in content:
            output_lines.append(f"--- Slide {slide['slide_number']} ---")
            if slide["title"]:
                output_lines.append(f"Title: {slide['title']}")
            if slide["content"]:
                output_lines.append("Content:")
                for content_item in slide["content"]:
                    output_lines.append(f"  - {content_item}")
            if slide["notes"]:
                output_lines.append(f"Notes: {slide['notes']}")
            output_lines.append("")

        return LoaderOutput(
            output_text="\n".join(output_lines),
            content=content,
            metadata=result_state.data.get("metadata", {}),
            slide_count=len(content),
            extracted_slides=content,
            data=result_state.data
        )

    except Exception as e:
        return LoaderOutput(
            output_text="",
            success=False,
            error_message=str(e)
        )
