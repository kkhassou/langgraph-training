import os
from typing import List, Dict, Any
from pptx import Presentation
import aiofiles

from src.nodes.base import BaseNode, NodeState, NodeInput, NodeOutput


class PowerPointIngestNode(BaseNode):
    """Node for ingesting and extracting text from PowerPoint files"""

    def __init__(self):
        super().__init__(
            name="ppt_ingest",
            description="Extract text content from PowerPoint presentations"
        )

    async def execute(self, state: NodeState) -> NodeState:
        """Execute PowerPoint text extraction"""
        try:
            file_path = state.data.get("file_path")
            if not file_path:
                raise ValueError("file_path is required in state.data")

            if not os.path.exists(file_path):
                raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

            # Extract text from PowerPoint
            extracted_text = await self._extract_text_from_ppt(file_path)

            # Update state
            state.data["extracted_text"] = extracted_text
            state.data["slide_count"] = len(extracted_text)
            state.messages.append(f"Extracted text from {len(extracted_text)} slides")
            state.metadata["node"] = self.name

            return state

        except Exception as e:
            state.data["error"] = str(e)
            state.metadata["error_node"] = self.name
            return state

    async def _extract_text_from_ppt(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint slides"""
        slides_text = []

        # Load presentation
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "notes": ""
            }

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    
                    # Try to identify title
                    is_title = False
                    try:
                        # Check if it's a placeholder and if it's a title
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                is_title = True
                    except Exception:
                        # If accessing placeholder_format fails, treat as regular content
                        pass
                    
                    if is_title:
                        slide_content["title"] = text_content
                    else:
                        slide_content["content"].append(text_content)

            # Extract notes
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide_content["notes"] = slide.notes_slide.notes_text_frame.text.strip()

            slides_text.append(slide_content)

        return slides_text


class PPTIngestInput(NodeInput):
    """Input model for PowerPoint ingest node"""
    file_path: str


class PPTIngestOutput(NodeOutput):
    """Output model for PowerPoint ingest node"""
    slide_count: int = 0
    extracted_slides: List[Dict[str, Any]] = []


async def ppt_ingest_handler(file_path: str) -> PPTIngestOutput:
    """Standalone handler for PowerPoint ingest API endpoint"""
    try:
        node = PowerPointIngestNode()
        state = NodeState()
        state.data["file_path"] = file_path

        result_state = await node.execute(state)

        if "error" in result_state.data:
            return PPTIngestOutput(
                output_text="",
                success=False,
                error_message=result_state.data["error"]
            )

        extracted_text = result_state.data.get("extracted_text", [])

        # Format output text
        output_lines = []
        for slide in extracted_text:
            output_lines.append(f"--- Slide {slide['slide_number']} ---")
            if slide["title"]:
                output_lines.append(f"Title: {slide['title']}")
            if slide["content"]:
                output_lines.append("Content:")
                for content in slide["content"]:
                    output_lines.append(f"  - {content}")
            if slide["notes"]:
                output_lines.append(f"Notes: {slide['notes']}")
            output_lines.append("")

        return PPTIngestOutput(
            output_text="\n".join(output_lines),
            slide_count=result_state.data.get("slide_count", 0),
            extracted_slides=extracted_text,
            data=result_state.data
        )

    except Exception as e:
        return PPTIngestOutput(
            output_text="",
            success=False,
            error_message=str(e)
        )