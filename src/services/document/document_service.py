"""Document Service - ドキュメント処理の統合サービス"""

from typing import List, Dict, Any
from pydantic import BaseModel
import logging
from pptx import Presentation

logger = logging.getLogger(__name__)


class SlideContent(BaseModel):
    """スライドコンテンツ"""
    slide_number: int
    title: str
    content: List[str]
    notes: str


class DocumentService:
    """ドキュメント処理サービス
    
    各種ドキュメントフォーマットの読み込み・解析を提供
    """
    
    @staticmethod
    async def extract_from_ppt(file_path: str) -> List[SlideContent]:
        """
        PowerPointファイルからテキストを抽出
        
        Args:
            file_path: PPTファイルのパス
        
        Returns:
            スライドコンテンツのリスト
        
        Example:
            >>> slides = await DocumentService.extract_from_ppt("presentation.pptx")
            >>> for slide in slides:
            ...     print(f"Slide {slide.slide_number}: {slide.title}")
        """
        logger.info(f"Extracting text from PowerPoint: {file_path}")
        
        slides_content = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "notes": ""
            }
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    
                    # Try to identify title
                    is_title = False
                    try:
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                is_title = True
                    except Exception:
                        pass
                    
                    if is_title:
                        slide_data["title"] = text_content
                    else:
                        slide_data["content"].append(text_content)
            
            # Extract notes
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
            
            slides_content.append(SlideContent(**slide_data))
        
        logger.info(f"Extracted {len(slides_content)} slides")
        return slides_content
    
    @staticmethod
    async def extract_from_pdf(file_path: str) -> List[Dict[str, Any]]:
        """
        PDFファイルからテキストを抽出（将来の実装）
        
        Args:
            file_path: PDFファイルのパス
        
        Returns:
            ページコンテンツのリスト
        """
        # TODO: Implement PDF extraction
        raise NotImplementedError("PDF extraction not yet implemented")
    
    @staticmethod
    async def extract_from_word(file_path: str) -> Dict[str, Any]:
        """
        Wordファイルからテキストを抽出（将来の実装）
        
        Args:
            file_path: Wordファイルのパス
        
        Returns:
            ドキュメントコンテンツ
        """
        # TODO: Implement Word extraction
        raise NotImplementedError("Word extraction not yet implemented")
    
    @staticmethod
    def format_slides_as_text(slides: List[SlideContent]) -> str:
        """
        スライドをテキスト形式にフォーマット
        
        Args:
            slides: スライドコンテンツのリスト
        
        Returns:
            フォーマットされたテキスト
        """
        output_lines = []
        
        for slide in slides:
            output_lines.append(f"--- Slide {slide.slide_number} ---")
            
            if slide.title:
                output_lines.append(f"Title: {slide.title}")
            
            if slide.content:
                output_lines.append("Content:")
                for content in slide.content:
                    output_lines.append(f"  - {content}")
            
            if slide.notes:
                output_lines.append(f"Notes: {slide.notes}")
            
            output_lines.append("")
        
        return "\n".join(output_lines)


# 便利な関数エイリアス
async def extract_ppt_text(file_path: str) -> List[SlideContent]:
    """PPTテキスト抽出のエイリアス"""
    return await DocumentService.extract_from_ppt(file_path)

